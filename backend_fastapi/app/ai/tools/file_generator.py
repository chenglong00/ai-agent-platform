"""Tools for generating files: PDF, DOCX, PPTX, HTML, and Markdown reports."""

import json
from pathlib import Path

from langchain.tools import tool
from pydantic import BaseModel, Field

_OUTPUT_DIR = Path("/tmp/deep_agent_outputs")


def _ensure_output_dir() -> Path:
    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return _OUTPUT_DIR


# ── Input schemas ────────────────────────────────────────────────────────────

class HtmlReportInput(BaseModel):
    filename: str = Field(description="Output filename without extension")
    title: str = Field(description="Document title shown in <h1> and <title>")
    html_body: str = Field(description="Full HTML markup for the body (everything that goes inside <body>)")
    extra_css: str = Field(default="", description="Optional additional CSS styles")


class SectionedDocInput(BaseModel):
    filename: str = Field(description="Output filename without extension")
    title: str = Field(description="Document title")
    sections: str = Field(
        description=(
            'JSON array of sections. Each item: {"heading": "Section Title", "content": "Paragraph text..."}. '
            'Use \\n inside content for line breaks.'
        )
    )


class PptxInput(BaseModel):
    filename: str = Field(description="Output filename without extension")
    title: str = Field(description="Presentation title (shown on title slide)")
    subtitle: str = Field(default="", description="Optional subtitle for the title slide")
    slides: str = Field(
        description=(
            'JSON array of slides. Each item: {"title": "Slide Title", "bullets": ["Point 1", "Point 2", ...]}. '
            'The first slide is auto-generated as a title slide; these are content slides.'
        )
    )


class MarkdownReportInput(BaseModel):
    filename: str = Field(description="Output filename without extension")
    content: str = Field(description="Full Markdown content to save")


# ── HTML ─────────────────────────────────────────────────────────────────────

_DEFAULT_CSS = """
body { font-family: Arial, Helvetica, sans-serif; max-width: 960px; margin: 48px auto; padding: 0 24px; color: #222; line-height: 1.6; }
h1 { font-size: 2em; color: #1a1a2e; border-bottom: 3px solid #4f8ef7; padding-bottom: 12px; margin-bottom: 24px; }
h2 { font-size: 1.4em; color: #16213e; margin-top: 32px; }
h3 { font-size: 1.1em; color: #0f3460; }
p  { margin: 10px 0; }
ul, ol { padding-left: 24px; }
li { margin: 6px 0; }
table { border-collapse: collapse; width: 100%; margin: 24px 0; }
th, td { border: 1px solid #ccc; padding: 10px 14px; text-align: left; }
th { background: #1a1a2e; color: #fff; }
tr:nth-child(even) { background: #f5f7fa; }
.highlight { background: #fff8e1; border-left: 5px solid #ffb300; padding: 14px 18px; margin: 24px 0; border-radius: 4px; }
.summary   { background: #e8f4fd; padding: 16px 20px; border-radius: 8px; margin: 24px 0; }
footer { margin-top: 48px; font-size: 0.85em; color: #888; border-top: 1px solid #ddd; padding-top: 12px; }
"""


@tool(args_schema=HtmlReportInput)
def create_html_report(filename: str, title: str, html_body: str, extra_css: str = "") -> str:
    """Create an HTML report file and return its path. Use for web-viewable reports."""
    out = _ensure_output_dir()
    path = out / f"{filename}.html"
    path.write_text(
        f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{title}</title>
  <style>
{_DEFAULT_CSS}
{extra_css}
  </style>
</head>
<body>
{html_body}
</body>
</html>""",
        encoding="utf-8",
    )
    return str(path)


# ── PDF ──────────────────────────────────────────────────────────────────────

@tool(args_schema=SectionedDocInput)
def create_pdf_report(filename: str, title: str, sections: str) -> str:
    """Create a PDF report from sections and return its path. Requires reportlab."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    except ImportError:
        return (
            "ERROR: 'reportlab' is not installed. "
            "Run: pip install reportlab  — then retry."
        )

    try:
        items: list[dict] = json.loads(sections)
    except json.JSONDecodeError as exc:
        return f"ERROR: sections must be valid JSON — {exc}"

    out = _ensure_output_dir()
    path = out / f"{filename}.pdf"

    doc = SimpleDocTemplate(str(path), pagesize=letter, leftMargin=inch, rightMargin=inch)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "DocTitle", parent=styles["Title"],
        fontSize=22, spaceAfter=18, textColor=colors.HexColor("#1a1a2e"),
    )
    h2_style = ParagraphStyle(
        "DocH2", parent=styles["Heading2"],
        fontSize=14, spaceBefore=16, spaceAfter=8, textColor=colors.HexColor("#16213e"),
    )
    body_style = ParagraphStyle(
        "DocBody", parent=styles["Normal"],
        fontSize=11, leading=17, spaceAfter=6,
    )

    story = [Paragraph(title, title_style), HRFlowable(width="100%", thickness=2), Spacer(1, 0.2 * inch)]

    for section in items:
        heading = str(section.get("heading") or "").strip()
        content = str(section.get("content") or "").strip()
        if heading:
            story.append(Paragraph(heading, h2_style))
        for line in content.splitlines():
            if line.strip():
                story.append(Paragraph(line.strip(), body_style))
        story.append(Spacer(1, 0.15 * inch))

    doc.build(story)
    return str(path)


# ── DOCX ─────────────────────────────────────────────────────────────────────

@tool(args_schema=SectionedDocInput)
def create_docx_document(filename: str, title: str, sections: str) -> str:
    """Create a Word (.docx) document from sections and return its path. Requires python-docx."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return (
            "ERROR: 'python-docx' is not installed. "
            "Run: pip install python-docx  — then retry."
        )

    try:
        items: list[dict] = json.loads(sections)
    except json.JSONDecodeError as exc:
        return f"ERROR: sections must be valid JSON — {exc}"

    doc = Document()

    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.runs[0]
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    for section in items:
        heading = str(section.get("heading") or "").strip()
        content = str(section.get("content") or "").strip()
        if heading:
            h = doc.add_heading(heading, level=2)
            h.runs[0].font.color.rgb = RGBColor(0x16, 0x21, 0x3E)
        if content:
            for line in content.splitlines():
                if line.strip():
                    p = doc.add_paragraph(line.strip())
                    p.runs[0].font.size = Pt(11)

    out = _ensure_output_dir()
    path = out / f"{filename}.docx"
    doc.save(str(path))
    return str(path)


# ── PPTX ─────────────────────────────────────────────────────────────────────

@tool(args_schema=PptxInput)
def create_pptx_presentation(filename: str, title: str, subtitle: str, slides: str) -> str:
    """Create a PowerPoint (.pptx) presentation and return its path. Requires python-pptx."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
    except ImportError:
        return (
            "ERROR: 'python-pptx' is not installed. "
            "Run: pip install python-pptx  — then retry."
        )

    try:
        slide_data: list[dict] = json.loads(slides)
    except json.JSONDecodeError as exc:
        return f"ERROR: slides must be valid JSON — {exc}"

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    ts = prs.slides.add_slide(title_layout)
    ts.shapes.title.text = title
    if subtitle and ts.placeholders[1]:
        ts.placeholders[1].text = subtitle

    # Content slides
    content_layout = prs.slide_layouts[1]
    for sd in slide_data:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = str(sd.get("title") or "")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        bullets: list[str] = sd.get("bullets") or []
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(bullet)
            p.level = 0

    out = _ensure_output_dir()
    path = out / f"{filename}.pptx"
    prs.save(str(path))
    return str(path)


# ── Markdown ─────────────────────────────────────────────────────────────────

@tool(args_schema=MarkdownReportInput)
def create_markdown_report(filename: str, content: str) -> str:
    """Save content as a Markdown (.md) file and return its path. Always available."""
    out = _ensure_output_dir()
    path = out / f"{filename}.md"
    path.write_text(content, encoding="utf-8")
    return str(path)
